#!/usr/bin/env python3
"""
LocalfileAgent.py — Scan files/directories and either summarise them or
chat with them interactively using a local Ollama model.

Modes
-----
  summarise (default)  — generate a summary for every file
  chat                 — load files into context, then ask questions in a REPL

Usage
-----
  # After `pip install -e .` the `localfileagent` command is on your PATH:
  localfileagent /path/to/directory
  localfileagent file1.txt file2.py --output summaries.md
  localfileagent /path/to/dir --ext .py .md --recursive
  localfileagent /path/to/dir --chat
  localfileagent /path/to/dir --chat --model gemma3
  localfileagent --gui                      # launch the graphical interface

  # Or run the scripts directly without installing:
  python LocalfileAgent.py /path/to/directory
  python LocalfileAgent.py --gui
"""

import argparse
import json
import socket
import sys
import time
import urllib.request
import urllib.error
from pathlib import Path
from typing import Iterator

# ── Configuration ─────────────────────────────────────────────────────────────

DEFAULT_MODEL     = "mistral"
DEFAULT_EMBED_MODEL = "nomic-embed-text"
OLLAMA_GENERATE   = "http://localhost:11434/api/generate"
OLLAMA_CHAT       = "http://localhost:11434/api/chat"
OLLAMA_TAGS       = "http://localhost:11434/api/tags"
OLLAMA_EMBED      = "http://localhost:11434/api/embed"
REQUEST_TIMEOUT   = 600       # seconds; generous ceiling for slow local generation
MAX_FILE_BYTES    = 200_000   # skip text files larger than ~200 KB
MAX_EXTRACT_CHARS = 400_000   # cap extracted text from binary docs
CONTEXT_FILE_CAP  = 20        # max files loaded into chat context
MAX_CLARIFY_ROUNDS = 3        # max clarifying questions per user turn

TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".java", ".c", ".cpp", ".h", ".cs", ".go", ".rs",
    ".rb", ".php", ".swift", ".kt", ".sh", ".bash",
    ".html", ".css", ".json", ".yaml", ".yml", ".toml",
    ".xml", ".csv", ".rst", ".sql",
    # Config / environment
    ".ini", ".cfg", ".conf", ".env",
    # IaC / DevOps
    ".tf", ".tfvars", ".hcl", ".ps1", ".bat", ".cmd",
    # Additional languages
    ".dart", ".scala", ".lua", ".ex", ".exs",
    ".r", ".zig",
    # Web frameworks
    ".vue", ".svelte",
    # Schema / IDL
    ".graphql", ".gql", ".proto",
}

BINARY_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    # Notebooks / email
    ".ipynb", ".eml",
    # OpenDocument
    ".odt", ".ods", ".odp",
    # Archives
    ".zip", ".7z", ".tar", ".tgz",
}

SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | BINARY_EXTENSIONS

SUMMARISE_SYSTEM = (
    "You are a concise technical assistant. "
    "Summarise the provided file content in 3-5 sentences. "
    "Focus on what the file does, its main components, and any noteworthy patterns. "
    "Be direct and factual. Do not repeat the filename."
)

GENERAL_SYSTEM = (
    "You are a helpful assistant. Answer the user's questions accurately and "
    "concisely, drawing on your own knowledge."
)

CHAT_SYSTEM_TEMPLATE = """\
You are a helpful assistant with access to the following {n} file(s).
Answer questions about their content accurately and concisely, and cite which
file you drew from when you use it. You may also answer questions the files
don't cover from your own knowledge — don't refuse just because the answer
isn't in the files.

{file_block}
"""

RAG_SYSTEM = (
    "You are a helpful assistant. The user's message may include context "
    "excerpts pulled from their local files, each labelled with its source file "
    "in square brackets. When an excerpt is relevant, use it and cite the source "
    "file. When the excerpts don't cover the question, answer normally from your "
    "own knowledge — do not refuse or say the answer isn't in the files."
)

_CLARIFY_PROBE = (
    "[META-TASK: Do not answer the user's previous message yet.]\n"
    "Rate your confidence (0–100) that you fully understand what is being asked "
    "given the conversation context.\n\n"
    "If confidence < 95, respond ONLY with:\n"
    "CONFIDENCE: <number>\n"
    "QUESTION: <one concise clarifying question>\n\n"
    "If confidence >= 95, respond ONLY with:\n"
    "CONFIDENCE: 95\n"
    "PROCEED"
)

# ── Ollama helpers ────────────────────────────────────────────────────────────

def _post(url: str, payload: dict, timeout: int = REQUEST_TIMEOUT) -> dict:
    data = json.dumps(payload).encode()
    req = urllib.request.Request(
        url, data=data,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            return json.loads(resp.read())
    except urllib.error.HTTPError as exc:
        body = ""
        try:
            body = exc.read().decode("utf-8", errors="replace").strip()
        except Exception:
            pass
        model = payload.get("model", "?")
        if exc.code == 404:
            raise ValueError(
                f"Ollama model '{model}' is not pulled.\n"
                f"Run:  ollama pull {model}"
            ) from exc
        raise ConnectionError(
            f"Ollama returned HTTP {exc.code} for model '{model}'."
            + (f"\n{body}" if body else "")
        ) from exc
    except urllib.error.URLError as exc:
        if isinstance(exc.reason, (TimeoutError, socket.timeout)):
            raise TimeoutError(
                f"Ollama did not respond within {timeout}s — "
                "the file may be too large for this model."
            ) from exc
        raise ConnectionError(
            "Cannot reach Ollama at http://localhost:11434.\n"
            "Make sure the local Ollama daemon is running ('ollama serve')."
        ) from exc


def query_ollama_generate(prompt: str, system: str, model: str) -> str:
    """Single-turn generation (used for summarisation)."""
    result = _post(OLLAMA_GENERATE, {
        "model": model,
        "prompt": prompt,
        "system": system,
        "stream": False,
    })
    return result.get("response", "").strip()


def query_ollama_chat(messages: list[dict], model: str) -> tuple[str, list[dict]]:
    """
    Multi-turn chat.  Returns (assistant_reply, updated_messages).
    The system message is baked into messages[0].
    """
    result = _post(OLLAMA_CHAT, {
        "model": model,
        "messages": messages,
        "stream": False,
    })
    assistant_msg = result.get("message", {})
    reply = assistant_msg.get("content", "").strip()
    return reply, messages + [{"role": "assistant", "content": reply}]


def stream_ollama_chat(messages: list[dict], model: str) -> Iterator[str]:
    """Yields token strings one at a time from the Ollama streaming chat API."""
    data = json.dumps({
        "model": model,
        "messages": messages,
        "stream": True,
    }).encode()
    req = urllib.request.Request(
        OLLAMA_CHAT, data=data,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=REQUEST_TIMEOUT) as resp:
            phase: str | None = None  # None | "thinking" | "content"
            for raw_line in resp:
                line = raw_line.strip()
                if not line:
                    continue
                try:
                    obj = json.loads(line)
                except json.JSONDecodeError:
                    continue
                msg = obj.get("message", {})
                thinking = msg.get("thinking", "")
                content = msg.get("content", "")
                if thinking:
                    if phase != "thinking":
                        yield "💭 " if phase is None else "\n\n💭 "
                        phase = "thinking"
                    yield thinking
                if content:
                    if phase == "thinking":
                        yield "\n\n"
                    phase = "content"
                    yield content
                if obj.get("done"):
                    break
    except urllib.error.URLError as exc:
        if isinstance(exc.reason, (TimeoutError, socket.timeout)):
            raise TimeoutError(
                "Ollama did not respond — the file may be too large for this model."
            ) from exc
        raise ConnectionError(
            "Cannot reach Ollama at http://localhost:11434.\n"
            "Make sure the local Ollama daemon is running ('ollama serve')."
        ) from exc


def embed_ollama(texts: list[str], model: str) -> list[list[float]]:
    """Embed a batch of strings via Ollama's /api/embed. Returns one vector per input."""
    result = _post(OLLAMA_EMBED, {"model": model, "input": texts})
    embeddings = result.get("embeddings")
    if not embeddings:
        raise ValueError(
            f"Ollama returned no embeddings — check that model '{model}' "
            f"supports embeddings (e.g. ollama pull nomic-embed-text)."
        )
    return embeddings


def check_ollama_available(model: str, embed_model: str | None = None) -> None:
    """Verify Ollama is reachable and the requested model(s) are pulled."""
    try:
        req = urllib.request.Request(OLLAMA_TAGS)
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json.loads(resp.read())
        available = [m["name"] for m in data["models"]]   # full names, tags kept
    except urllib.error.URLError:
        print("✗  Ollama is not running.  Start it with:  ollama serve", file=sys.stderr)
        sys.exit(1)
    except (json.JSONDecodeError, TypeError, KeyError) as exc:
        print(f"✗  Unexpected response from Ollama at {OLLAMA_TAGS}: {exc}", file=sys.stderr)
        sys.exit(1)

    # A bare name (e.g. "mistral") matches any pulled tag; a tagged request
    # (e.g. "mistral:7b") must match exactly.
    available_bases = {n.split(":")[0] for n in available}
    wanted = [model] + ([embed_model] if embed_model else [])
    missing = []
    for m in wanted:
        present = (m in available) if ":" in m else (m.split(":")[0] in available_bases)
        if not present:
            missing.append(m)
    if missing:
        print(
            f"⚠  Model(s) not found locally: {', '.join(missing)}\n"
            f"   Available: {', '.join(available) or 'none'}\n"
            f"   Pull with:  {'; '.join(f'ollama pull {m}' for m in missing)}\n",
            file=sys.stderr,
        )
        sys.exit(1)


def run_clarification_check(messages: list[dict], model: str) -> tuple[int, str | None]:
    """Non-destructively probe the LLM for clarity on the most recent user message.

    Returns (confidence, question) where question is None when the LLM is >=95%
    confident. Never raises — returns (95, None) on any error so chat is never
    blocked by a failed probe.
    """
    probe = list(messages) + [{"role": "user", "content": _CLARIFY_PROBE}]
    try:
        reply, _ = query_ollama_chat(probe, model)
    except Exception:
        return 95, None

    confidence = 95
    question = None
    for line in reply.splitlines():
        s = line.strip()
        upper = s.upper()
        if upper.startswith("CONFIDENCE:"):
            raw = s.split(":", 1)[1].strip()
            digits = "".join(c for c in raw if c.isdigit())[:3]
            if digits:
                try:
                    confidence = int(digits)
                except ValueError:
                    pass
        elif upper.startswith("QUESTION:"):
            question = s.split(":", 1)[1].strip()

    if confidence >= 95 or "PROCEED" in reply.upper() or not question:
        return 95, None
    return confidence, question


# ── File collection ───────────────────────────────────────────────────────────

def collect_files(paths: list[str], extensions: set[str], recursive: bool) -> list[Path]:
    collected: list[Path] = []
    for raw in paths:
        p = Path(raw).expanduser().resolve()
        if not p.exists():
            print(f"⚠  Skipping (not found): {p}", file=sys.stderr)
            continue
        if p.is_file():
            if p.suffix.lower() in extensions:
                collected.append(p)
            else:
                print(f"⚠  Skipping unsupported type: {p.name}", file=sys.stderr)
        elif p.is_dir():
            glob = "**/*" if recursive else "*"
            for child in sorted(p.glob(glob)):
                if child.is_file() and child.suffix.lower() in extensions:
                    collected.append(child)

    seen: set[Path] = set()
    unique: list[Path] = []
    for f in collected:
        if f not in seen:
            seen.add(f)
            unique.append(f)
    return unique


def read_file_safe(path: Path) -> str | None:
    """Return file text, or None if it should be skipped."""
    try:
        size = path.stat().st_size
        if size == 0:
            return None
        ext = path.suffix.lower()
        if ext in BINARY_EXTENSIONS:
            text = _extract_binary(path, ext)
            if text is None:
                return None
            return text[:MAX_EXTRACT_CHARS]
        if size > MAX_FILE_BYTES:
            return None
        return path.read_text(encoding="utf-8", errors="replace")
    except OSError:
        return None


def _extract_binary(path: Path, ext: str) -> str | None:
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        if ext == ".docx":
            return _extract_docx(path)
        if ext == ".xlsx":
            return _extract_xlsx(path)
        if ext == ".pptx":
            return _extract_pptx(path)
        if ext == ".xls":
            return _extract_xls(path)
        if ext in (".doc", ".ppt"):
            return _extract_via_word_powerpoint(path, ext)
        if ext == ".ipynb":
            return _extract_ipynb(path)
        if ext == ".eml":
            return _extract_eml(path)
        if ext in (".odt", ".ods", ".odp"):
            return _extract_odf(path, ext)
        if ext in (".zip", ".7z", ".tar", ".tgz"):
            return _extract_archive(path, ext)
    except Exception as exc:
        print(f"\n⚠  Failed to extract {path.name}: {exc}", file=sys.stderr)
        return None
    return None


def _missing(pkg: str, ext: str) -> None:
    print(
        f"\n⚠  Cannot read {ext} files — install '{pkg}':  pip install {pkg}",
        file=sys.stderr,
    )


def _extract_pdf(path: Path) -> str | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        _missing("pypdf", ".pdf")
        return None
    reader = PdfReader(str(path))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_docx(path: Path) -> str | None:
    try:
        import docx  # python-docx
    except ImportError:
        _missing("python-docx", ".docx")
        return None
    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str | None:
    try:
        from openpyxl import load_workbook
    except ImportError:
        _missing("openpyxl", ".xlsx")
        return None
    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            parts.append("\t".join("" if v is None else str(v) for v in row))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str | None:
    try:
        from pptx import Presentation
    except ImportError:
        _missing("python-pptx", ".pptx")
        return None
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_xls(path: Path) -> str | None:
    try:
        import xlrd
    except ImportError:
        _missing("xlrd", ".xls")
        return None
    book = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"[Sheet: {sheet.name}]")
        for r in range(sheet.nrows):
            parts.append("\t".join(str(sheet.cell_value(r, c)) for c in range(sheet.ncols)))
    return "\n".join(parts)


def _extract_via_word_powerpoint(path: Path, ext: str) -> str | None:
    """Legacy .doc/.ppt — uses Microsoft Office via COM (Windows only)."""
    try:
        import win32com.client  # pywin32
    except ImportError:
        _missing("pywin32", ext)
        return None
    import pythoncom
    pythoncom.CoInitialize()
    try:
        if ext == ".doc":
            app = win32com.client.Dispatch("Word.Application")
            app.Visible = False
            try:
                doc = app.Documents.Open(str(path), ReadOnly=True)
                try:
                    return doc.Content.Text
                finally:
                    doc.Close(SaveChanges=False)
            finally:
                app.Quit()
        else:  # .ppt
            app = win32com.client.Dispatch("PowerPoint.Application")
            try:
                pres = app.Presentations.Open(str(path), WithWindow=False, ReadOnly=True)
                try:
                    parts = []
                    for i, slide in enumerate(pres.Slides, 1):
                        parts.append(f"[Slide {i}]")
                        for shape in slide.Shapes:
                            if shape.HasTextFrame and shape.TextFrame.HasText:
                                parts.append(shape.TextFrame.TextRange.Text)
                    return "\n".join(parts)
                finally:
                    pres.Close()
            finally:
                app.Quit()
    finally:
        pythoncom.CoUninitialize()


def _extract_ipynb(path: Path) -> str | None:
    try:
        nb = json.loads(path.read_bytes())
    except Exception:
        return None
    parts: list[str] = []
    for cell in nb.get("cells", []):
        cell_type = cell.get("cell_type", "")
        source = "".join(cell.get("source", []))
        if not source.strip():
            continue
        if cell_type == "code":
            parts.append(f"```python\n{source}\n```")
        else:
            parts.append(source)
        for output in cell.get("outputs", []):
            text = output.get("text") or output.get("data", {}).get("text/plain")
            if isinstance(text, list):
                text = "".join(text)
            if text and text.strip():
                parts.append(f"[Output]\n{text}")
    return "\n\n".join(parts) or None


def _extract_eml(path: Path) -> str | None:
    import email as _email
    import email.policy
    msg = _email.message_from_bytes(path.read_bytes(), policy=_email.policy.compat32)
    parts: list[str] = []
    for header in ("Subject", "From", "To", "Date"):
        val = msg.get(header)
        if val:
            parts.append(f"{header}: {val}")
    parts.append("")
    for part in msg.walk():
        if part.get_content_type() == "text/plain":
            payload = part.get_payload(decode=True)
            if payload:
                parts.append(payload.decode("utf-8", errors="replace"))
    return "\n".join(parts) or None


def _extract_odf(path: Path, ext: str) -> str | None:
    try:
        from odf.opendocument import load as odf_load
        from odf import text as odf_text
        from odf import teletype
    except ImportError:
        _missing("odfpy", ext)
        return None
    doc = odf_load(str(path))
    paragraphs = doc.getElementsByType(odf_text.P)
    return "\n".join(teletype.extractText(p) for p in paragraphs) or None


def _extract_archive(path: Path, ext: str) -> str | None:
    parts: list[str] = []
    total = 0

    def _fits(name: str) -> bool:
        return Path(name).suffix.lower() in TEXT_EXTENSIONS

    def _add(name: str, content: str) -> bool:
        nonlocal total
        parts.append(f"[{name}]\n{content}")
        total += len(content)
        return total >= MAX_EXTRACT_CHARS

    if ext == ".zip":
        import zipfile
        if not zipfile.is_zipfile(path):
            return None
        with zipfile.ZipFile(path) as zf:
            for info in zf.infolist():
                if info.is_dir() or not _fits(info.filename):
                    continue
                try:
                    with zf.open(info) as f:
                        if _add(info.filename, f.read().decode("utf-8", errors="replace")):
                            break
                except Exception:
                    continue

    elif ext == ".7z":
        try:
            import py7zr
        except ImportError:
            _missing("py7zr", ".7z")
            return None
        with py7zr.SevenZipFile(path, mode="r") as z:
            names = [n for n in z.getnames() if _fits(n)]
            if names:
                for name, bio in z.read(names).items():
                    if _add(name, bio.read().decode("utf-8", errors="replace")):
                        break

    elif ext in (".tar", ".tgz"):
        import tarfile
        try:
            with tarfile.open(path, "r:*") as tf:
                for member in tf.getmembers():
                    if not member.isfile() or not _fits(member.name):
                        continue
                    try:
                        f = tf.extractfile(member)
                        if f is None:
                            continue
                        if _add(member.name, f.read().decode("utf-8", errors="replace")):
                            break
                    except Exception:
                        continue
        except tarfile.TarError:
            return None

    return "\n\n".join(parts) or None


# ── Summarise mode ────────────────────────────────────────────────────────────

def run_summarise(files: list[Path], model: str, output: str | None) -> None:
    print(f"📂  Found {len(files)} file(s). Summarising with '{model}'…\n")

    results: list[tuple[Path, str]] = []
    for i, path in enumerate(files, 1):
        print(f"[{i}/{len(files)}] {path.name}", end=" … ", flush=True)
        elapsed = 0.0
        content = read_file_safe(path)
        if content is None:
            try:
                size = path.stat().st_size
            except OSError:
                summary = "(skipped — file no longer accessible)"
            else:
                summary = "(empty file)" if size == 0 else f"(skipped — too large: {size:,} bytes)"
        else:
            try:
                t0 = time.monotonic()
                summary = query_ollama_generate(
                    f"File: {path.name}\n\n{content}",
                    SUMMARISE_SYSTEM,
                    model,
                )
                elapsed = time.monotonic() - t0
            except TimeoutError as exc:
                summary = f"(skipped — {exc})"
                elapsed = 0.0
            except ConnectionError as exc:
                print(f"\n✗  {exc}", file=sys.stderr)
                sys.exit(1)
        results.append((path, summary))
        suffix = f" ({elapsed:.1f}s)" if elapsed else ""
        print(f"done{suffix}")

    output_path = Path(output) if output else None
    use_md = output_path and output_path.suffix.lower() == ".md"
    text = _fmt_markdown(results) if use_md else _fmt_plain(results)

    if output_path:
        output_path.write_text(text, encoding="utf-8")
        print(f"\n✅  Summaries written to {output_path}")
    else:
        print(f"\n{'─'*60}\n{text}")


def _fmt_markdown(results: list[tuple[Path, str]]) -> str:
    lines = ["# File Summaries\n"]
    for path, summary in results:
        lines += [f"## `{path.name}`", f"**Path:** `{path}`\n", summary, ""]
    return "\n".join(lines)


def _fmt_plain(results: list[tuple[Path, str]]) -> str:
    sep = "─" * 60
    parts = []
    for path, summary in results:
        parts += [sep, f"FILE : {path}", sep, summary, ""]
    return "\n".join(parts)


# ── Chat mode ─────────────────────────────────────────────────────────────────

def build_file_block(files: list[Path]) -> tuple[str, list[str]]:
    """Return (formatted file block for the system prompt, list of skipped names)."""
    parts: list[str] = []
    skipped: list[str] = []
    for path in files:
        content = read_file_safe(path)
        if content is None:
            skipped.append(path.name)
            continue
        parts.append(f"### {path.name}\nPath: {path}\n\n{content}")
    return "\n\n---\n\n".join(parts), skipped


def run_chat(files: list[Path], model: str, *, embed_model: str = DEFAULT_EMBED_MODEL,
             top_k: int = 5, use_rag: bool = True) -> None:
    if len(files) > CONTEXT_FILE_CAP:
        print(
            f"⚠  {len(files)} files found — only the first {CONTEXT_FILE_CAP} will be "
            f"used (limit).\n   Use --ext or --recursive to narrow the selection.\n"
        )
        files = files[:CONTEXT_FILE_CAP]

    if not files:
        # No files loaded — plain general-assistant chat.
        messages = [{"role": "system", "content": GENERAL_SYSTEM}]
        print(
            "💬  Chat mode — no files loaded, general chat mode. Add files anytime "
            "for grounded answers.\n"
            f"    Commands:  /clear  /quit\n{'─'*60}"
        )
        _chat_repl(messages, model, index=None, embed_model=embed_model, top_k=top_k)
        return

    index = None
    if use_rag:
        try:
            print(f"📂  Indexing {len(files)} file(s) with '{embed_model}'…", end=" ", flush=True)
            from rag import build_index
            index = build_index(files, embed_model)
            if len(index) == 0:
                print("\n⚠  No content indexed — falling back to full-context mode.")
                index = None
            else:
                print(f"done  ({len(index)} chunks)")
        except (ImportError, ValueError) as exc:
            # numpy missing, or the embed model returned no usable vectors.
            print(f"\n⚠  {exc}\n   Falling back to full-context mode.")
            index = None
        except (ConnectionError, TimeoutError) as exc:
            print(f"\n✗  {exc}", file=sys.stderr)
            sys.exit(1)

    if index is not None:
        messages: list[dict] = [{"role": "system", "content": RAG_SYSTEM}]
        print(
            f"\n💬  Chat mode (RAG) — ask anything about the indexed files.\n"
            f"    Commands:  /clear  /quit\n{'─'*60}"
        )
    else:
        # Fallback: original full-context behavior.
        print(f"📂  Loading {len(files)} file(s) into context…", end=" ", flush=True)
        file_block, skipped = build_file_block(files)
        if not file_block.strip():
            print("\n✗  No readable content found in the selected files.", file=sys.stderr)
            sys.exit(1)
        loaded = len(files) - len(skipped)
        print(f"done  ({loaded} loaded, {len(skipped)} skipped)")
        if skipped:
            print(f"   Skipped (empty/too large): {', '.join(skipped)}")
        messages = [{"role": "system",
                     "content": CHAT_SYSTEM_TEMPLATE.format(n=loaded, file_block=file_block)}]
        print(
            f"\n💬  Chat mode — ask anything about the loaded files.\n"
            f"    Commands:  /clear  /quit\n{'─'*60}"
        )

    _chat_repl(messages, model, index=index, embed_model=embed_model, top_k=top_k)


def _chat_repl(messages: list[dict], model: str, *, index, embed_model: str,
               top_k: int) -> None:
    """The shared REPL loop. `index` is a VectorIndex (RAG) or None (general /
    full-context chat); when None the plain history is sent as-is."""
    while True:
        try:
            user_input = input("\nYou: ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\n\nBye!")
            break
        if not user_input:
            continue
        if user_input.lower() in ("/quit", "/exit", "/q"):
            print("Bye!")
            break
        if user_input.lower() == "/clear":
            messages[:] = [messages[0]]
            print("🗑  Conversation history cleared.")
            continue

        # Persisted history keeps the plain user text; only the outgoing copy
        # gets the retrieved context injected into the latest turn.
        messages.append({"role": "user", "content": user_input})
        orig_idx = len(messages) - 1   # position of the original question

        # Clarification loop — ask at most MAX_CLARIFY_ROUNDS questions before
        # proceeding; the Q&A exchange is kept in the persisted history.
        for _ in range(MAX_CLARIFY_ROUNDS):
            confidence, question = run_clarification_check(messages, model)
            if confidence >= 95 or question is None:
                break
            print(f"\n{model} [clarifying]: {question}")
            messages.append({"role": "assistant", "content": question})
            try:
                clarify_answer = input("You: ").strip()
            except (EOFError, KeyboardInterrupt):
                print("\n\nBye!")
                return
            if not clarify_answer:
                break
            messages.append({"role": "user", "content": clarify_answer})

        if index is not None:
            from rag import retrieve, build_rag_prompt
            try:
                chunks = retrieve(index, user_input, embed_model, top_k)
            except (ConnectionError, TimeoutError) as exc:
                print(f"\n✗  {exc}", file=sys.stderr)
                sys.exit(1)
            # Inject RAG context into the original question, not the last message,
            # so that clarification answers remain intact in the API call.
            api_messages = (
                messages[:orig_idx]
                + [{"role": "user", "content": build_rag_prompt(chunks, user_input)}]
                + messages[orig_idx + 1:]
            )
        else:
            api_messages = messages

        print(f"\n{model}: ", end="", flush=True)
        try:
            t0 = time.monotonic()
            reply, _ = query_ollama_chat(api_messages, model)
            elapsed = time.monotonic() - t0
        except (ConnectionError, TimeoutError) as exc:
            print(f"\n✗  {exc}", file=sys.stderr)
            sys.exit(1)
        messages.append({"role": "assistant", "content": reply})
        print(reply)
        print(f"  ⏱  {elapsed:.1f}s")


# ── CLI ───────────────────────────────────────────────────────────────────────

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Summarise files or chat with them using a local Ollama model.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    p.add_argument("paths", nargs="*", help="Files or directories to scan (omit when using --gui).")
    p.add_argument(
        "--gui",
        action="store_true",
        help="Launch the graphical interface instead of the CLI.",
    )
    p.add_argument(
        "--chat", "-c",
        action="store_true",
        help="Interactive Q&A mode instead of summarisation.",
    )
    p.add_argument(
        "--model", "-m",
        default=DEFAULT_MODEL,
        help=f"Ollama model to use (default: {DEFAULT_MODEL}).",
    )
    p.add_argument(
        "--embed-model",
        default=DEFAULT_EMBED_MODEL,
        help=f"Embedding model for RAG (default: {DEFAULT_EMBED_MODEL}).",
    )
    p.add_argument(
        "--top-k",
        type=int,
        default=5,
        help="(chat) Number of context chunks to retrieve per question (default: 5).",
    )
    p.add_argument(
        "--no-rag",
        action="store_true",
        help="(chat) Disable RAG; load full file contents into context instead.",
    )
    p.add_argument(
        "--output", "-o",
        help="(summarise mode) Write output to this file; .md extension → Markdown.",
    )
    p.add_argument(
        "--ext",
        nargs="+",
        metavar="EXT",
        help="File extensions to include, e.g. --ext .py .md",
    )
    p.add_argument(
        "--recursive", "-r",
        action="store_true",
        help="Recurse into subdirectories.",
    )
    p.add_argument(
        "--no-check",
        action="store_true",
        help="Skip the Ollama availability check.",
    )
    return p


def _force_utf8_output(streams=None) -> None:
    """Make stdout/stderr tolerate the emoji/box-drawing glyphs we print, even
    when the console's locale encoding (e.g. cp1252 on Windows) cannot — without
    this a redirected or piped run dies with UnicodeEncodeError on the first ✓."""
    for stream in (streams if streams is not None else (sys.stdout, sys.stderr)):
        reconfigure = getattr(stream, "reconfigure", None)
        if reconfigure is None:
            continue
        try:
            reconfigure(encoding="utf-8")
        except (ValueError, OSError):
            pass


def main() -> None:
    _force_utf8_output()
    parser = build_parser()
    args = parser.parse_args()

    if args.gui:
        try:
            import gui
        except ImportError as exc:
            print(f"✗  The GUI needs PySide6 — install it:  pip install PySide6\n   ({exc})",
                  file=sys.stderr)
            sys.exit(1)
        gui.main()
        return

    # Chat with no paths → general-assistant chat (no files). Summarise still
    # needs files, so it keeps requiring at least one path.
    if not args.paths and not args.chat:
        parser.error("the following arguments are required: paths (or pass --gui / --chat)")

    extensions = (
        {e if e.startswith(".") else f".{e}" for e in args.ext}
        if args.ext else SUPPORTED_EXTENSIONS
    )

    if not args.no_check:
        embed_model = None if args.no_rag else args.embed_model
        check_ollama_available(args.model, embed_model)

    files = collect_files(args.paths, extensions, args.recursive) if args.paths else []
    if not files and not args.chat:
        print("No matching files found.", file=sys.stderr)
        sys.exit(0)

    if args.chat:
        run_chat(files, args.model, embed_model=args.embed_model,
                 top_k=args.top_k, use_rag=not args.no_rag)
    else:
        run_summarise(files, args.model, args.output)


if __name__ == "__main__":
    main()
